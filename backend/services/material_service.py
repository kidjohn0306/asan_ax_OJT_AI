"""
Google Drive 교육자료(education_materials 폴더) 스캔·캐싱.

폴더 구조: {DRIVE_EDUCATION_MATERIALS_FOLDER_ID}/ 아래 common, team1, team2, team3 등
category명과 동일한 하위 폴더 — admin_service.TEAM_KEY_MAP이 team_code를 이 category명으로
변환하므로(T1→team1 등) 폴더명을 그대로 캐시 key로 사용한다.

파일 id + modifiedTime을 이전 스캔 결과(repositories.material_repo)와 비교해 변경분만
재다운로드·재추출한다 — 매 문제 생성마다 Drive를 다시 훑거나 텍스트를 재추출하면
비용(Drive API 호출 + 추출 시간)이 쌓이기 때문에, 변경 없는 파일은 캐시된 텍스트를 그대로 재사용한다.
"""
import io
import logging
import os
from datetime import datetime, timezone

from services.drive_service import DriveService

EDUCATION_ROOT_FOLDER_ID_ENV = "DRIVE_EDUCATION_MATERIALS_FOLDER_ID"

_SUPPORTED_MIME_TYPES = {
    "application/pdf",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
}

_MAX_CHARS_PER_FILE = 20000  # 파일 1개당 추출 텍스트 상한 (프롬프트 폭주 방지)
_MAX_TOTAL_CHARS = 40000  # category 전체 합산 상한 — Sheets 셀 한도(5만자) 안에 여유를 두고 유지


def _max_file_size_bytes() -> int:
    # 대용량/압축폭탄성 파일로 인한 메모리 고갈 방지. 실제 교육자료(임베디드 영상·고해상도 이미지 포함 PPTX)가
    # 기본값보다 큰 경우를 대비해 환경변수로 조정 가능하게 함.
    return int(os.getenv("DRIVE_MATERIAL_MAX_FILE_SIZE_MB", "25")) * 1024 * 1024


def _root_folder_id() -> str:
    return os.getenv(EDUCATION_ROOT_FOLDER_ID_ENV, "")


def categories_for_team(team_code: str) -> list:
    """team_code에 해당하는 material 카테고리 목록(공통 + 팀별). team/admin_service 양쪽에서 공유."""
    from services.admin_service import TEAM_KEY_MAP

    category = TEAM_KEY_MAP.get(team_code, team_code)
    return ["common", category] if category != "common" else ["common"]


def _extract_text(data: bytes, mime_type: str) -> str:
    if mime_type == "application/pdf":
        from pypdf import PdfReader
        reader = PdfReader(io.BytesIO(data))
        text = "\n".join(page.extract_text() or "" for page in reader.pages)
    elif mime_type == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        from pptx import Presentation
        prs = Presentation(io.BytesIO(data))
        lines = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    lines.append(shape.text_frame.text)
        text = "\n".join(lines)
    else:
        return ""

    text = text.strip()
    if len(text) > _MAX_CHARS_PER_FILE:
        text = text[:_MAX_CHARS_PER_FILE] + "\n[이하 생략 — 파일이 너무 깁니다]"
    return text


def list_remote_files(category: str) -> list:
    """category 하위폴더의 지원 형식(PDF/PPTX) 파일 목록.
    루트 폴더 미설정이거나 하위폴더가 없으면 빈 리스트를 반환한다 (Drive 미연동 환경 대응)."""
    root_id = _root_folder_id()
    if not root_id:
        return []
    drive = DriveService()
    folder_id = drive.find_child_folder(root_id, category)
    if not folder_id:
        return []
    files = drive.list_files(folder_id, page_size=100)
    return [f for f in files if f.get("mimeType") in _SUPPORTED_MIME_TYPES]


def check_new_materials(category: str, manifest: dict | None = None) -> dict:
    """캐시된 매니페스트와 Drive 목록을 비교해 새로 추가되었거나 변경된 파일을 찾는다.
    Drive 조회 실패 시(미연동·권한 없음 등) 예외를 삼키고 '새 파일 없음'으로 취급해
    문제 생성/관리자 화면이 이 기능 때문에 깨지지 않도록 한다.
    manifest를 이미 조회해둔 호출자는 그대로 넘겨서 Sheets 재조회(쿼터 소모)를 피한다."""
    from repositories import material_repo

    try:
        remote_files = list_remote_files(category)
    except Exception as e:
        logging.warning(f"교육자료 목록 조회 실패 (category={category}): {e}")
        return {"category": category, "has_new": False, "new_files": [], "cached_file_count": 0, "scanned_at": ""}

    if manifest is None:
        manifest = material_repo.get_manifest(category) or {}
    cached_modified = {f["id"]: f.get("modifiedTime") for f in manifest.get("files", [])}

    new_files = [f for f in remote_files if cached_modified.get(f["id"]) != f.get("modifiedTime")]

    return {
        "category": category,
        "has_new": bool(new_files),
        "new_files": [{"id": f["id"], "name": f.get("name", "")} for f in new_files],
        "cached_file_count": len(manifest.get("files", [])),
        "scanned_at": manifest.get("scanned_at", ""),
    }


def scan_materials(category: str) -> dict:
    """새/변경된 파일만 다운로드+텍스트추출하고, 기존 캐시와 병합해 저장한다."""
    from repositories import material_repo

    remote_files = list_remote_files(category)
    manifest = material_repo.get_manifest(category) or {"category": category, "files": []}
    cached_by_id = {f["id"]: f for f in manifest.get("files", [])}

    max_size = _max_file_size_bytes()
    drive = DriveService()
    updated_files = []
    skipped = []
    for f in remote_files:
        cached = cached_by_id.get(f["id"])
        # modifiedTime이 같아도 이전 스캔에서 추출에 실패(용량 초과 등)했던 파일은 재시도한다 —
        # 그렇지 않으면 admin이 DRIVE_MATERIAL_MAX_FILE_SIZE_MB를 올려도 파일이 영영 스캔되지 않는다.
        if cached and cached.get("modifiedTime") == f.get("modifiedTime") and cached.get("extracted", True):
            updated_files.append(cached)
            continue
        extracted = True
        try:
            file_size = int(f.get("size") or 0)
            if file_size > max_size:
                raise ValueError(f"파일 크기 초과 ({file_size} bytes > {max_size} bytes)")
            data = drive.download_bytes(f["id"])
            text = _extract_text(data, f.get("mimeType", ""))
        except Exception as e:
            logging.warning(f"교육자료 다운로드/추출 실패 (file={f.get('name')}): {e}")
            text = cached.get("text", "") if cached else ""  # 실패 시 이전에 캐시된 텍스트라도 유지
            extracted = False  # 다음 스캔 때 재시도 대상으로 표시
            skipped.append({"name": f.get("name", ""), "reason": str(e)})
        updated_files.append({
            "id": f["id"],
            "name": f.get("name", ""),
            "mimeType": f.get("mimeType", ""),
            "modifiedTime": f.get("modifiedTime", ""),
            "text": text,
            "extracted": extracted,
        })

    # category 전체 텍스트 합산이 저장소 셀 한도를 넘지 않도록, 파일 순서대로 예산을 소진하며 자른다.
    # (combined_text를 별도로 저장하지 않고 files[].text에서 매번 계산하므로, 여기서 잘라두면
    #  저장되는 텍스트 총량 자체가 줄어 Sheets 셀 크기 문제와 저장 중복을 함께 없앤다.)
    budget = _MAX_TOTAL_CHARS
    capped_files = []
    for uf in updated_files:
        text = uf.get("text", "")
        if budget <= 0:
            text = ""
        elif len(text) > budget:
            text = text[:budget] + "\n[이하 생략 — 저장 용량 제한]"
        budget -= len(text)
        capped_files.append({**uf, "text": text})

    new_manifest = {
        "category": category,
        "files": capped_files,
        "scanned_at": datetime.now(timezone.utc).isoformat(),
    }
    material_repo.save_manifest(category, new_manifest)
    return {**new_manifest, "skipped": skipped}


def list_cached_materials(team_code: str | None = None) -> dict:
    """캐시된 교육자료 전체 목록을 category별로 반환한다 (텍스트 본문은 응답 크기 때문에 제외).
    team_code가 없으면 등록된 모든 팀 + 공통 카테고리를 합쳐서 보여준다."""
    from repositories import material_repo, team_repo

    if team_code:
        categories = categories_for_team(team_code)
        labels = {"common": "공통"}
        if len(categories) > 1:
            team_name = next(
                (t.get("team_name", team_code) for t in team_repo.list_teams() if t.get("team_code") == team_code),
                team_code,
            )
            labels[categories[-1]] = team_name
    else:
        team_rows = team_repo.list_teams()
        categories = ["common"]
        labels = {"common": "공통"}
        for t in team_rows:
            code = t.get("team_code", "")
            if not code:
                continue
            cat = categories_for_team(code)[-1]
            if cat not in categories:
                categories.append(cat)
            labels.setdefault(cat, t.get("team_name", cat))

    result = {}
    for cat in categories:
        manifest = material_repo.get_manifest(cat) or {}
        cached = manifest.get("files", [])
        cached_ids = {f.get("id") for f in cached}
        # Drive에 새로 올라왔거나 변경된 파일을 함께 표시 — 아직 스캔 전이라 캐시에는 없거나 구버전으로 남아있다.
        new_check = check_new_materials(cat, manifest=manifest)
        new_ids = {f["id"] for f in new_check.get("new_files", [])}

        files = [
            {
                "id": f.get("id", ""),
                "name": f.get("name", ""),
                "mimeType": f.get("mimeType", ""),
                "modifiedTime": f.get("modifiedTime", ""),
                "status": "new" if f.get("id") in new_ids else ("synced" if f.get("extracted", True) else "failed"),
            }
            for f in cached
        ]
        files += [
            {"id": nf["id"], "name": nf.get("name", ""), "mimeType": "", "modifiedTime": "", "status": "new"}
            for nf in new_check.get("new_files", [])
            if nf["id"] not in cached_ids
        ]

        result[cat] = {
            "category": cat,
            "label": labels.get(cat, cat),
            "files": files,
            "scanned_at": manifest.get("scanned_at", ""),
            "has_new": bool(new_check.get("new_files")),
        }
    return {"categories": result}


def get_cached_text(category: str, selected_ids: set | None = None) -> str:
    """캐시된 파일별 텍스트를 결합해 반환. combined_text를 별도 저장하지 않고
    매번 files[]에서 계산해 중복 저장을 피한다 (비용 미미한 문자열 결합일 뿐).
    selected_ids가 주어지면 관리자가 화면에서 선택 해제한 파일은 제외한다."""
    from repositories import material_repo
    manifest = material_repo.get_manifest(category)
    if not manifest:
        return ""
    files = manifest.get("files", [])
    if selected_ids is not None:
        files = [f for f in files if f.get("id") in selected_ids]
    return "\n\n".join(f"[{f['name']}]\n{f['text']}" for f in files if f.get("text"))


def get_material_text_for_team(team_code: str, selected_ids: set | None = None) -> str:
    """공통(common) + 팀별 캐시 텍스트를 합쳐 반환 — AI 문제 생성 시 자동 소스로 사용된다.
    selected_ids가 None이면(기본값) 캐시된 파일 전체를 포함한다."""
    parts = [text for cat in categories_for_team(team_code) if (text := get_cached_text(cat, selected_ids))]
    return "\n\n".join(parts)
